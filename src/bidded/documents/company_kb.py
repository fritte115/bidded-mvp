from __future__ import annotations

import csv
import hashlib
import re
from collections.abc import Mapping, Sequence
from dataclasses import dataclass
from enum import StrEnum
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, Protocol
from uuid import UUID

import fitz
from pydantic import BaseModel, ConfigDict, Field

DEMO_TENANT_KEY = "demo"
DEFAULT_COMPANY_KB_BUCKET = "company-knowledge"
COMPANY_KB_CHUNKING_STRATEGY = "company_kb_text_unit_char_budget_v1"
DEFAULT_MAX_CHUNK_CHARS = 1800
_DRAFT_ATTACHMENT_TYPES = {
    "certificate",
    "cv",
    "reference_case",
    "policy_document",
    "pricing_document",
    "other",
}


class CompanyKbError(RuntimeError):
    """Raised when a company knowledge-base document cannot be handled."""


class CompanyKbRegistrationError(CompanyKbError):
    """Raised when a draft-oriented company KB PDF cannot be registered."""


class CompanyKbDocumentType(StrEnum):
    CERTIFICATION = "certification"
    CASE_STUDY = "case_study"
    CV_PROFILE = "cv_profile"
    CAPABILITY_STATEMENT = "capability_statement"
    POLICY_PROCESS = "policy_process"
    FINANCIAL_PRICING = "financial_pricing"
    LEGAL_INSURANCE = "legal_insurance"


_DOCUMENT_TYPE_BY_DRAFT_ATTACHMENT = {
    "certificate": CompanyKbDocumentType.CERTIFICATION,
    "cv": CompanyKbDocumentType.CV_PROFILE,
    "reference_case": CompanyKbDocumentType.CASE_STUDY,
    "policy_document": CompanyKbDocumentType.POLICY_PROCESS,
    "pricing_document": CompanyKbDocumentType.FINANCIAL_PRICING,
    "other": CompanyKbDocumentType.CAPABILITY_STATEMENT,
}


_ALLOWED_EXTENSIONS: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ),
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".csv": "text/csv",
    ".txt": "text/plain",
    ".md": "text/markdown",
}

_CATEGORY_BY_DOCUMENT_TYPE: dict[CompanyKbDocumentType, str] = {
    CompanyKbDocumentType.CERTIFICATION: "certification",
    CompanyKbDocumentType.CASE_STUDY: "reference",
    CompanyKbDocumentType.CV_PROFILE: "cv_summary",
    CompanyKbDocumentType.CAPABILITY_STATEMENT: "capability",
    CompanyKbDocumentType.POLICY_PROCESS: "policy",
    CompanyKbDocumentType.FINANCIAL_PRICING: "economics",
    CompanyKbDocumentType.LEGAL_INSURANCE: "legal",
}


class SupabaseCompanyKbQuery(Protocol):
    def select(self, columns: str) -> SupabaseCompanyKbQuery: ...

    def eq(self, column: str, value: object) -> SupabaseCompanyKbQuery: ...

    def upsert(
        self,
        payload: dict[str, Any] | list[dict[str, Any]],
        *,
        on_conflict: str | None = None,
    ) -> SupabaseCompanyKbQuery: ...

    def insert(self, payload: list[dict[str, Any]]) -> SupabaseCompanyKbQuery: ...

    def update(self, payload: dict[str, Any]) -> SupabaseCompanyKbQuery: ...

    def delete(self) -> SupabaseCompanyKbQuery: ...

    def execute(self) -> Any: ...


class SupabaseCompanyKbStorageBucket(Protocol):
    def upload(
        self,
        path: str,
        file: bytes,
        *,
        file_options: dict[str, str] | None = None,
    ) -> Any: ...

    def download(self, path: str) -> bytes: ...

    def remove(self, paths: list[str]) -> Any: ...


class SupabaseCompanyKbStorage(Protocol):
    def from_(self, bucket_name: str) -> SupabaseCompanyKbStorageBucket: ...


class SupabaseCompanyKbClient(Protocol):
    storage: SupabaseCompanyKbStorage

    def table(self, table_name: str) -> SupabaseCompanyKbQuery: ...


@dataclass(frozen=True)
class CompanyKbUploadFile:
    filename: str
    content: bytes
    content_type: str
    kb_document_type: CompanyKbDocumentType | str


@dataclass(frozen=True)
class CompanyKbRegistrationResult:
    company_id: UUID
    document_id: UUID
    storage_path: str
    checksum_sha256: str
    content_type: str
    original_filename: str
    kb_document_type: CompanyKbDocumentType


@dataclass(frozen=True)
class CompanyKbPdfRegistrationResult:
    company_id: str
    document_id: str
    storage_path: str
    checksum_sha256: str
    content_type: str
    original_filename: str
    attachment_type: str


@dataclass(frozen=True)
class CompanyKbIngestionResult:
    document_id: UUID
    company_id: UUID
    chunk_count: int
    evidence_count: int
    extraction_status: str
    warnings: tuple[str, ...] = ()


@dataclass(frozen=True)
class CompanyKbDocumentSummary:
    document_id: UUID
    company_id: UUID
    original_filename: str
    storage_path: str
    content_type: str
    parse_status: str
    kb_document_type: CompanyKbDocumentType
    extraction_status: str
    evidence_count: int
    warnings: tuple[str, ...]


@dataclass(frozen=True)
class ExtractedTextUnit:
    unit_number: int
    text: str
    metadata: dict[str, Any]


@dataclass(frozen=True)
class CompanyKbChunk:
    document_id: UUID
    page_start: int
    page_end: int
    chunk_index: int
    text: str
    metadata: dict[str, Any]

    def to_payload(self) -> dict[str, Any]:
        return {
            "tenant_key": DEMO_TENANT_KEY,
            "document_id": str(self.document_id),
            "page_start": self.page_start,
            "page_end": self.page_end,
            "chunk_index": self.chunk_index,
            "text": self.text,
            "metadata": self.metadata,
        }


class CompanyKbChunkForExtraction(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    chunk_id: UUID
    chunk_index: int
    page_start: int
    page_end: int
    text: str = Field(min_length=1)
    metadata: dict[str, Any] = Field(default_factory=dict)


class EvidenceExcerptRef(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    chunk_id: UUID
    excerpt: str = Field(min_length=1)


class ExtractedCompanyKbFact(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    fact_type: str = Field(min_length=1)
    category: str = Field(min_length=1)
    claim: str = Field(min_length=1)
    normalized_meaning: str = Field(min_length=1)
    evidence_ref: EvidenceExcerptRef
    confidence: float = Field(ge=0, le=1)
    metadata: dict[str, Any] = Field(default_factory=dict)


class ExtractedCompanyKbFacts(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    facts: tuple[ExtractedCompanyKbFact, ...] = ()


class CompanyKbExtractionRequest(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    company_id: UUID
    document_id: UUID
    kb_document_type: CompanyKbDocumentType
    original_filename: str
    chunks: tuple[CompanyKbChunkForExtraction, ...]


class CompanyKbFactExtractor(Protocol):
    def extract(
        self, request: CompanyKbExtractionRequest
    ) -> ExtractedCompanyKbFacts: ...


def register_company_kb_documents(
    client: SupabaseCompanyKbClient,
    *,
    company_id: UUID | str,
    bucket_name: str = DEFAULT_COMPANY_KB_BUCKET,
    files: Sequence[CompanyKbUploadFile],
    created_via: str = "bidded_api",
) -> list[CompanyKbRegistrationResult]:
    """Upload and register company knowledge-base files as company_profile docs."""

    normalized_company_id = _normalize_uuid(company_id, "company_id")
    results: list[CompanyKbRegistrationResult] = []
    for file in files:
        kb_document_type = _coerce_document_type(file.kb_document_type)
        content_type = _content_type_for_file(file.filename, file.content_type)
        if not file.content:
            raise CompanyKbError(f"Company KB file is empty: {file.filename}")

        checksum_sha256 = hashlib.sha256(file.content).hexdigest()
        storage_path = _storage_path(
            company_id=normalized_company_id,
            checksum_sha256=checksum_sha256,
            original_filename=file.filename,
        )
        client.storage.from_(bucket_name).upload(
            storage_path,
            file.content,
            file_options={"content-type": content_type, "upsert": "true"},
        )
        response = (
            client.table("documents")
            .upsert(
                {
                    "tenant_key": DEMO_TENANT_KEY,
                    "tender_id": None,
                    "company_id": str(normalized_company_id),
                    "storage_path": storage_path,
                    "checksum_sha256": checksum_sha256,
                    "content_type": content_type,
                    "document_role": "company_profile",
                    "parse_status": "pending",
                    "original_filename": file.filename,
                    "metadata": {
                        "registered_via": created_via,
                        "source_label": file.filename,
                        "kb_document_type": kb_document_type.value,
                        "extraction_status": "pending",
                        "warnings": [],
                    },
                },
                on_conflict="storage_path",
            )
            .execute()
        )
        row = _first_row(response, "documents")
        results.append(
            CompanyKbRegistrationResult(
                company_id=normalized_company_id,
                document_id=_normalize_uuid(row.get("id"), "documents.id"),
                storage_path=storage_path,
                checksum_sha256=checksum_sha256,
                content_type=content_type,
                original_filename=file.filename,
                kb_document_type=kb_document_type,
            )
        )

    return results


def register_company_kb_pdf(
    client: SupabaseCompanyKbClient,
    *,
    pdf_path: Path,
    bucket_name: str,
    company_id: UUID | str,
    source_label: str,
    attachment_type: str = "other",
    created_via: str = "bidded_cli",
) -> CompanyKbPdfRegistrationResult:
    """Register an approved PDF using the draft attachment vocabulary."""

    normalized_company_id = _normalize_uuid(company_id, "company_id")
    normalized_attachment_type = _normalize_draft_attachment_type(attachment_type)
    kb_document_type = _DOCUMENT_TYPE_BY_DRAFT_ATTACHMENT[normalized_attachment_type]
    path = Path(pdf_path)
    pdf_bytes = _read_valid_draft_pdf(path)
    checksum_sha256 = hashlib.sha256(pdf_bytes).hexdigest()
    storage_path = _draft_storage_path(
        company_id=normalized_company_id,
        checksum_sha256=checksum_sha256,
        original_filename=path.name,
    )
    client.storage.from_(bucket_name).upload(
        storage_path,
        pdf_bytes,
        file_options={"content-type": "application/pdf", "upsert": "true"},
    )
    response = (
        client.table("documents")
        .upsert(
            {
                "tenant_key": DEMO_TENANT_KEY,
                "tender_id": None,
                "company_id": str(normalized_company_id),
                "storage_path": storage_path,
                "checksum_sha256": checksum_sha256,
                "content_type": "application/pdf",
                "document_role": "company_profile",
                "parse_status": "pending",
                "original_filename": path.name,
                "metadata": {
                    "registered_via": created_via,
                    "source_label": source_label.strip() or path.name,
                    "kb_document_type": kb_document_type.value,
                    "kb_attachment_type": normalized_attachment_type,
                    "approved_for_bid_drafts": True,
                    "extraction_status": "pending",
                    "warnings": [],
                },
            },
            on_conflict="storage_path",
        )
        .execute()
    )
    row = _first_row(response, "documents")
    return CompanyKbPdfRegistrationResult(
        company_id=str(normalized_company_id),
        document_id=str(row["id"]),
        storage_path=storage_path,
        checksum_sha256=checksum_sha256,
        content_type="application/pdf",
        original_filename=path.name,
        attachment_type=normalized_attachment_type,
    )


def ingest_company_kb_document(
    client: SupabaseCompanyKbClient,
    *,
    document_id: UUID | str,
    bucket_name: str = DEFAULT_COMPANY_KB_BUCKET,
    max_chunk_chars: int = DEFAULT_MAX_CHUNK_CHARS,
    extractor: CompanyKbFactExtractor | None = None,
) -> CompanyKbIngestionResult:
    """Parse one company KB document and materialize active company evidence."""

    normalized_document_id = _normalize_uuid(document_id, "document_id")
    if max_chunk_chars <= 0:
        raise CompanyKbError("max_chunk_chars must be greater than zero.")

    row = _fetch_company_kb_document(client, normalized_document_id)
    company_id = _normalize_uuid(row.get("company_id"), "documents.company_id")
    metadata = dict(_mapping(row.get("metadata")))
    kb_document_type = _coerce_document_type(metadata.get("kb_document_type"))
    source_metadata = _source_metadata(row, kb_document_type=kb_document_type)

    try:
        _update_document_status(
            client,
            document_id=normalized_document_id,
            parse_status="parsing",
            metadata={
                **metadata,
                "extraction_status": "parsing",
                "parser": {
                    "status": "parsing",
                    "parser": "company_kb",
                    "chunking_strategy": COMPANY_KB_CHUNKING_STRATEGY,
                },
            },
        )
        file_bytes = client.storage.from_(bucket_name).download(
            str(row["storage_path"])
        )
        units = _extract_text_units(
            file_bytes,
            content_type=str(row.get("content_type") or ""),
            filename=str(row.get("original_filename") or ""),
        )
        chunks = build_company_kb_chunks(
            document_id=normalized_document_id,
            units=units,
            source_metadata=source_metadata,
            max_chunk_chars=max_chunk_chars,
        )
        if not chunks:
            raise CompanyKbError("No extractable text found in company KB document.")

        _replace_document_chunks(client, normalized_document_id, chunks)
        persisted_chunks = _fetch_chunks_for_document(client, normalized_document_id)
        facts, extraction_status, warnings = _extract_facts_or_fallback(
            extractor or RuleBasedCompanyKbFactExtractor(),
            request=CompanyKbExtractionRequest(
                company_id=company_id,
                document_id=normalized_document_id,
                kb_document_type=kb_document_type,
                original_filename=str(row.get("original_filename") or ""),
                chunks=tuple(persisted_chunks),
            ),
        )
        evidence_rows = _company_kb_evidence_payloads(
            company_id=company_id,
            document_id=normalized_document_id,
            kb_document_type=kb_document_type,
            original_filename=str(row.get("original_filename") or ""),
            facts=facts,
            chunks_by_id={chunk.chunk_id: chunk for chunk in persisted_chunks},
            extraction_status=extraction_status,
        )
        rows_returned = _upsert_company_kb_evidence(client, evidence_rows)
        parsed_metadata = {
            **metadata,
            "source_label": source_metadata["source_label"],
            "kb_document_type": kb_document_type.value,
            "extraction_status": extraction_status,
            "warnings": list(warnings),
            "parser": {
                "status": "parsed",
                "parser": "company_kb",
                "chunk_count": len(chunks),
                "chunking_strategy": COMPANY_KB_CHUNKING_STRATEGY,
            },
        }
        _update_document_status(
            client,
            document_id=normalized_document_id,
            parse_status="parsed",
            metadata=parsed_metadata,
        )
        return CompanyKbIngestionResult(
            document_id=normalized_document_id,
            company_id=company_id,
            chunk_count=len(chunks),
            evidence_count=rows_returned or len(evidence_rows),
            extraction_status=extraction_status,
            warnings=warnings,
        )
    except CompanyKbError as exc:
        _update_document_status(
            client,
            document_id=normalized_document_id,
            parse_status="parser_failed",
            metadata={
                **metadata,
                "extraction_status": "failed",
                "warnings": [str(exc)],
                "parser": {
                    "status": "parser_failed",
                    "parser": "company_kb",
                    "error_message": str(exc),
                },
            },
        )
        raise


def build_company_kb_chunks(
    *,
    document_id: UUID,
    units: Sequence[ExtractedTextUnit],
    source_metadata: Mapping[str, Any],
    max_chunk_chars: int = DEFAULT_MAX_CHUNK_CHARS,
) -> list[CompanyKbChunk]:
    chunks: list[CompanyKbChunk] = []
    source_label = str(source_metadata.get("source_label") or "company KB document")
    for unit in units:
        for segment_index, text in enumerate(
            _split_text(unit.text, max_chunk_chars=max_chunk_chars)
        ):
            chunks.append(
                CompanyKbChunk(
                    document_id=document_id,
                    page_start=unit.unit_number,
                    page_end=unit.unit_number,
                    chunk_index=len(chunks),
                    text=text,
                    metadata={
                        **dict(unit.metadata),
                        "source_label": source_label,
                        "source_document_id": str(document_id),
                        "source_original_filename": source_metadata.get(
                            "original_filename"
                        ),
                        "kb_document_type": source_metadata.get("kb_document_type"),
                        "parser": "company_kb",
                        "chunking_strategy": COMPANY_KB_CHUNKING_STRATEGY,
                        "max_chunk_chars": max_chunk_chars,
                        "unit_segment_index": segment_index,
                    },
                )
            )
    return chunks


def list_company_kb_documents(
    client: SupabaseCompanyKbClient,
    *,
    company_id: UUID | str,
) -> list[CompanyKbDocumentSummary]:
    normalized_company_id = _normalize_uuid(company_id, "company_id")
    doc_rows = _response_rows(
        client.table("documents")
        .select(
            "id,tenant_key,company_id,storage_path,content_type,document_role,"
            "parse_status,original_filename,metadata"
        )
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("company_id", str(normalized_company_id))
        .eq("document_role", "company_profile")
        .execute()
    )
    evidence_rows = _response_rows(
        client.table("evidence_items")
        .select("id,tenant_key,source_type,company_id,document_id")
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("source_type", "company_profile")
        .eq("company_id", str(normalized_company_id))
        .execute()
    )
    counts: dict[str, int] = {}
    for row in evidence_rows:
        document_id = row.get("document_id")
        if document_id is not None:
            counts[str(document_id)] = counts.get(str(document_id), 0) + 1

    summaries = [
        _document_summary(row, evidence_count=counts.get(str(row.get("id")), 0))
        for row in doc_rows
    ]
    return sorted(summaries, key=lambda item: item.original_filename.casefold())


def list_company_kb_evidence(
    client: SupabaseCompanyKbClient,
    *,
    company_id: UUID | str,
    document_id: UUID | str,
) -> list[dict[str, Any]]:
    normalized_company_id = _normalize_uuid(company_id, "company_id")
    normalized_document_id = _normalize_uuid(document_id, "document_id")
    rows = _response_rows(
        client.table("evidence_items")
        .select("*")
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("source_type", "company_profile")
        .eq("company_id", str(normalized_company_id))
        .eq("document_id", str(normalized_document_id))
        .execute()
    )
    return sorted((dict(row) for row in rows), key=lambda row: str(row["evidence_key"]))


def delete_company_kb_document(
    client: SupabaseCompanyKbClient,
    *,
    company_id: UUID | str,
    document_id: UUID | str,
    bucket_name: str = DEFAULT_COMPANY_KB_BUCKET,
) -> None:
    normalized_company_id = _normalize_uuid(company_id, "company_id")
    normalized_document_id = _normalize_uuid(document_id, "document_id")
    row = _fetch_company_kb_document(client, normalized_document_id)
    if (
        _normalize_uuid(row.get("company_id"), "documents.company_id")
        != normalized_company_id
    ):
        raise CompanyKbError(
            f"Company KB document does not belong to company {company_id}."
        )

    storage_path = str(row.get("storage_path") or "")
    if storage_path:
        client.storage.from_(bucket_name).remove([storage_path])
    (
        client.table("documents")
        .delete()
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("id", str(normalized_document_id))
        .eq("company_id", str(normalized_company_id))
        .execute()
    )


class RuleBasedCompanyKbFactExtractor:
    """Deterministic schema extractor used when no live LLM extractor is injected."""

    def extract(self, request: CompanyKbExtractionRequest) -> ExtractedCompanyKbFacts:
        if not request.chunks:
            return ExtractedCompanyKbFacts()
        chunk = request.chunks[0]
        text = chunk.text
        category = _CATEGORY_BY_DOCUMENT_TYPE[request.kb_document_type]
        if request.kb_document_type is CompanyKbDocumentType.CV_PROFILE:
            return self._cv_fact(request, chunk)

        claim_text = _summary_for_document_type(request.kb_document_type, text)
        return ExtractedCompanyKbFacts(
            facts=(
                ExtractedCompanyKbFact(
                    fact_type=request.kb_document_type.value,
                    category=category,
                    claim=claim_text,
                    normalized_meaning=claim_text,
                    evidence_ref=EvidenceExcerptRef(
                        chunk_id=chunk.chunk_id,
                        excerpt=_excerpt(text),
                    ),
                    confidence=0.72,
                    metadata={"extraction_method": "rule_based"},
                ),
            )
        )

    def _cv_fact(
        self,
        request: CompanyKbExtractionRequest,
        chunk: CompanyKbChunkForExtraction,
    ) -> ExtractedCompanyKbFacts:
        text = _anonymize_cv_text(chunk.text)
        claim = _summary_for_document_type(request.kb_document_type, text)
        return ExtractedCompanyKbFacts(
            facts=(
                ExtractedCompanyKbFact(
                    fact_type="cv_profile",
                    category="cv_summary",
                    claim=claim,
                    normalized_meaning=claim,
                    evidence_ref=EvidenceExcerptRef(
                        chunk_id=chunk.chunk_id,
                        excerpt=_excerpt(text),
                    ),
                    confidence=0.72,
                    metadata={"extraction_method": "rule_based", "anonymized": True},
                ),
            )
        )


def _extract_facts_or_fallback(
    extractor: CompanyKbFactExtractor,
    *,
    request: CompanyKbExtractionRequest,
) -> tuple[tuple[ExtractedCompanyKbFact, ...], str, tuple[str, ...]]:
    try:
        extracted = extractor.extract(request)
        facts = tuple(extracted.facts)
        _validate_fact_refs(
            facts,
            request.chunks,
            kb_document_type=request.kb_document_type,
        )
        if facts:
            return facts, "extracted", ()
    except Exception as exc:  # noqa: BLE001 - fallback is the product behavior
        return _fallback_facts(request), "fallback", (str(exc),)

    return _fallback_facts(request), "fallback", ("No structured facts extracted.",)


def _fallback_facts(
    request: CompanyKbExtractionRequest,
) -> tuple[ExtractedCompanyKbFact, ...]:
    facts: list[ExtractedCompanyKbFact] = []
    for chunk in request.chunks[:5]:
        text = (
            _anonymize_cv_text(chunk.text)
            if request.kb_document_type is CompanyKbDocumentType.CV_PROFILE
            else chunk.text
        )
        category = _CATEGORY_BY_DOCUMENT_TYPE[request.kb_document_type]
        claim = _summary_for_document_type(request.kb_document_type, text)
        facts.append(
            ExtractedCompanyKbFact(
                fact_type=request.kb_document_type.value,
                category=category,
                claim=claim,
                normalized_meaning=claim,
                evidence_ref=EvidenceExcerptRef(
                    chunk_id=chunk.chunk_id,
                    excerpt=_excerpt(text),
                ),
                confidence=0.55,
                metadata={
                    "fallback": True,
                    "anonymized": request.kb_document_type
                    is CompanyKbDocumentType.CV_PROFILE,
                },
            )
        )
    return tuple(facts)


def _company_kb_evidence_payloads(
    *,
    company_id: UUID,
    document_id: UUID,
    kb_document_type: CompanyKbDocumentType,
    original_filename: str,
    facts: Sequence[ExtractedCompanyKbFact],
    chunks_by_id: Mapping[UUID, CompanyKbChunkForExtraction],
    extraction_status: str,
) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for index, fact in enumerate(facts):
        chunk = chunks_by_id[fact.evidence_ref.chunk_id]
        fact_key = f"{kb_document_type.value}-{str(document_id)[:8]}-{index + 1}"
        rows.append(
            {
                "tenant_key": DEMO_TENANT_KEY,
                "evidence_key": f"COMPANY-KB-{_slug(fact_key)}",
                "source_type": "company_profile",
                "excerpt": fact.evidence_ref.excerpt,
                "normalized_meaning": fact.normalized_meaning,
                "category": fact.category,
                "confidence": fact.confidence,
                "source_metadata": {
                    "source_label": original_filename,
                    "original_filename": original_filename,
                    "kb_document_type": kb_document_type.value,
                },
                "document_id": str(document_id),
                "chunk_id": str(chunk.chunk_id),
                "page_start": chunk.page_start,
                "page_end": chunk.page_end,
                "company_id": str(company_id),
                "field_path": (
                    f"knowledge_base.{kb_document_type.value}.{document_id}.facts[{index}]"
                ),
                "metadata": {
                    **fact.metadata,
                    "claim": fact.claim,
                    "kb_document_type": kb_document_type.value,
                    "extraction_status": extraction_status,
                    "chunk_index": chunk.chunk_index,
                },
            }
        )
    return rows


def _upsert_company_kb_evidence(
    client: SupabaseCompanyKbClient,
    rows: Sequence[dict[str, Any]],
) -> int:
    if not rows:
        return 0
    response = (
        client.table("evidence_items")
        .upsert(list(rows), on_conflict="tenant_key,evidence_key")
        .execute()
    )
    data = getattr(response, "data", [])
    return len(data) if isinstance(data, list) else 0


def _validate_fact_refs(
    facts: Sequence[ExtractedCompanyKbFact],
    chunks: Sequence[CompanyKbChunkForExtraction],
    *,
    kb_document_type: CompanyKbDocumentType,
) -> None:
    chunks_by_id = {chunk.chunk_id: chunk for chunk in chunks}
    for fact in facts:
        chunk = chunks_by_id.get(fact.evidence_ref.chunk_id)
        if chunk is None:
            raise CompanyKbError(
                "Extracted company KB fact references an unknown chunk."
            )
        comparable_text = (
            _anonymize_cv_text(chunk.text)
            if kb_document_type is CompanyKbDocumentType.CV_PROFILE
            else chunk.text
        )
        if fact.evidence_ref.excerpt.casefold() not in comparable_text.casefold():
            raise CompanyKbError(
                "Extracted company KB fact excerpt does not resolve to its chunk."
            )


def _extract_text_units(
    file_bytes: bytes,
    *,
    content_type: str,
    filename: str,
) -> list[ExtractedTextUnit]:
    suffix = Path(filename).suffix.lower()
    if content_type == "application/pdf" or suffix == ".pdf":
        return _extract_pdf_units(file_bytes)
    if suffix == ".docx":
        return _extract_docx_units(file_bytes)
    if suffix == ".pptx":
        return _extract_pptx_units(file_bytes)
    if suffix == ".xlsx":
        return _extract_xlsx_units(file_bytes)
    if suffix == ".csv":
        return _extract_csv_units(file_bytes)
    if suffix in {".txt", ".md"}:
        return _extract_text_file_units(file_bytes, location_type="document")
    raise CompanyKbError(f"Unsupported company KB file type: {filename}")


def _extract_pdf_units(file_bytes: bytes) -> list[ExtractedTextUnit]:
    try:
        document = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise CompanyKbError(f"PDF extraction failed: {exc}") from exc
    try:
        units: list[ExtractedTextUnit] = []
        for page in document:
            text = _repair_pdf_text_extraction_artifacts(
                _normalize_text(page.get_text("text"))
            )
            if not text:
                continue
            units.append(
                ExtractedTextUnit(
                    unit_number=page.number + 1,
                    text=text,
                    metadata={
                        "location_type": "page",
                        "page_numbers": [page.number + 1],
                    },
                )
            )
        return units
    finally:
        document.close()


def _extract_docx_units(file_bytes: bytes) -> list[ExtractedTextUnit]:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - dependency contract covers this
        raise CompanyKbError("DOCX parsing requires python-docx.") from exc
    try:
        document = Document(BytesIO(file_bytes))
    except Exception as exc:
        raise CompanyKbError(f"DOCX extraction failed: {exc}") from exc
    lines = [
        _normalize_text(p.text) for p in document.paragraphs if _normalize_text(p.text)
    ]
    for table in document.tables:
        for row in table.rows:
            values = [_normalize_text(cell.text) for cell in row.cells]
            line = " | ".join(value for value in values if value)
            if line:
                lines.append(line)
    return _text_lines_to_units(lines, location_type="paragraph")


def _extract_pptx_units(file_bytes: bytes) -> list[ExtractedTextUnit]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise CompanyKbError("PPTX parsing requires python-pptx.") from exc
    try:
        presentation = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        raise CompanyKbError(f"PPTX extraction failed: {exc}") from exc
    units: list[ExtractedTextUnit] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = _normalize_text(getattr(shape, "text", ""))
            if text:
                parts.append(text)
        if parts:
            units.append(
                ExtractedTextUnit(
                    unit_number=index,
                    text="\n".join(parts),
                    metadata={"location_type": "slide", "slide_number": index},
                )
            )
    return units


def _extract_xlsx_units(file_bytes: bytes) -> list[ExtractedTextUnit]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise CompanyKbError("XLSX parsing requires openpyxl.") from exc
    try:
        workbook = load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as exc:
        raise CompanyKbError(f"XLSX extraction failed: {exc}") from exc
    units: list[ExtractedTextUnit] = []
    try:
        for index, sheet in enumerate(workbook.worksheets, start=1):
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                units.append(
                    ExtractedTextUnit(
                        unit_number=index,
                        text="\n".join(rows),
                        metadata={
                            "location_type": "sheet",
                            "sheet_name": sheet.title,
                        },
                    )
                )
    finally:
        workbook.close()
    return units


def _extract_csv_units(file_bytes: bytes) -> list[ExtractedTextUnit]:
    text = _decode_text(file_bytes)
    reader = csv.reader(StringIO(text))
    lines = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
    return _text_lines_to_units([line for line in lines if line], location_type="csv")


def _extract_text_file_units(
    file_bytes: bytes,
    *,
    location_type: str,
) -> list[ExtractedTextUnit]:
    text = _normalize_text(_decode_text(file_bytes))
    if not text:
        return []
    return [
        ExtractedTextUnit(
            unit_number=1,
            text=text,
            metadata={"location_type": location_type},
        )
    ]


def _text_lines_to_units(
    lines: Sequence[str], *, location_type: str
) -> list[ExtractedTextUnit]:
    text = _normalize_text("\n".join(lines))
    if not text:
        return []
    return [
        ExtractedTextUnit(
            unit_number=1,
            text=text,
            metadata={"location_type": location_type},
        )
    ]


def _replace_document_chunks(
    client: SupabaseCompanyKbClient,
    document_id: UUID,
    chunks: Sequence[CompanyKbChunk],
) -> None:
    (
        client.table("evidence_items")
        .delete()
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("source_type", "company_profile")
        .eq("document_id", str(document_id))
        .execute()
    )
    (
        client.table("document_chunks")
        .delete()
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("document_id", str(document_id))
        .execute()
    )
    (
        client.table("document_chunks")
        .insert([chunk.to_payload() for chunk in chunks])
        .execute()
    )


def _fetch_company_kb_document(
    client: SupabaseCompanyKbClient,
    document_id: UUID,
) -> dict[str, Any]:
    rows = _response_rows(
        client.table("documents")
        .select(
            "id,tenant_key,tender_id,company_id,storage_path,checksum_sha256,"
            "content_type,document_role,parse_status,original_filename,metadata"
        )
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("id", str(document_id))
        .execute()
    )
    if not rows:
        raise CompanyKbError(f"Company KB document does not exist: {document_id}")
    row = dict(rows[0])
    if row.get("document_role") != "company_profile":
        raise CompanyKbError(
            f"Document is not a company_profile document: {document_id}"
        )
    return row


def _fetch_chunks_for_document(
    client: SupabaseCompanyKbClient,
    document_id: UUID,
) -> list[CompanyKbChunkForExtraction]:
    rows = _response_rows(
        client.table("document_chunks")
        .select("id,document_id,page_start,page_end,chunk_index,text,metadata")
        .eq("document_id", str(document_id))
        .eq("tenant_key", DEMO_TENANT_KEY)
        .execute()
    )
    chunks = [
        CompanyKbChunkForExtraction(
            chunk_id=_normalize_uuid(row.get("id"), "document_chunks.id"),
            chunk_index=int(row.get("chunk_index") or 0),
            page_start=int(row.get("page_start") or 1),
            page_end=int(row.get("page_end") or row.get("page_start") or 1),
            text=str(row.get("text") or ""),
            metadata=dict(_mapping(row.get("metadata"))),
        )
        for row in rows
    ]
    return sorted(chunks, key=lambda chunk: chunk.chunk_index)


def _update_document_status(
    client: SupabaseCompanyKbClient,
    *,
    document_id: UUID,
    parse_status: str,
    metadata: Mapping[str, Any],
) -> None:
    (
        client.table("documents")
        .update({"parse_status": parse_status, "metadata": dict(metadata)})
        .eq("tenant_key", DEMO_TENANT_KEY)
        .eq("id", str(document_id))
        .execute()
    )


def _document_summary(
    row: Mapping[str, Any],
    *,
    evidence_count: int,
) -> CompanyKbDocumentSummary:
    metadata = _mapping(row.get("metadata"))
    return CompanyKbDocumentSummary(
        document_id=_normalize_uuid(row.get("id"), "documents.id"),
        company_id=_normalize_uuid(row.get("company_id"), "documents.company_id"),
        original_filename=str(row.get("original_filename") or ""),
        storage_path=str(row.get("storage_path") or ""),
        content_type=str(row.get("content_type") or ""),
        parse_status=str(row.get("parse_status") or "pending"),
        kb_document_type=_coerce_document_type(metadata.get("kb_document_type")),
        extraction_status=str(metadata.get("extraction_status") or "pending"),
        evidence_count=evidence_count,
        warnings=tuple(
            str(item) for item in metadata.get("warnings", []) if str(item).strip()
        ),
    )


def _source_metadata(
    row: Mapping[str, Any],
    *,
    kb_document_type: CompanyKbDocumentType,
) -> dict[str, Any]:
    metadata = _mapping(row.get("metadata"))
    source_label = str(
        metadata.get("source_label")
        or row.get("original_filename")
        or "company KB document"
    )
    return {
        "source_label": source_label,
        "original_filename": row.get("original_filename"),
        "storage_path": row.get("storage_path"),
        "kb_document_type": kb_document_type.value,
    }


def _first_row(response: Any, table_name: str) -> Mapping[str, Any]:
    rows = _response_rows(response)
    if rows:
        return rows[0]
    raise CompanyKbError(f"Supabase {table_name} write did not return a row.")


def _response_rows(response: Any) -> list[Mapping[str, Any]]:
    data = getattr(response, "data", None)
    if not isinstance(data, list):
        raise CompanyKbError("Supabase query did not return row data.")
    return [row for row in data if isinstance(row, Mapping)]


def _content_type_for_file(filename: str, content_type: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in _ALLOWED_EXTENSIONS:
        raise CompanyKbError(f"Unsupported company KB file type: {filename}")
    expected = _ALLOWED_EXTENSIONS[suffix]
    if content_type and content_type not in {
        expected,
        "application/octet-stream",
        "text/plain",
        "text/markdown",
        "text/csv",
    }:
        raise CompanyKbError(f"Unsupported company KB file type: {filename}")
    return expected


def _coerce_document_type(
    value: CompanyKbDocumentType | str | object,
) -> CompanyKbDocumentType:
    try:
        return CompanyKbDocumentType(str(value))
    except ValueError as exc:
        raise CompanyKbError(f"Unsupported company KB document type: {value}") from exc


def _storage_path(
    *,
    company_id: UUID,
    checksum_sha256: str,
    original_filename: str,
) -> str:
    return (
        f"demo/company-knowledge/{company_id}/"
        f"{checksum_sha256[:12]}-{_safe_filename(original_filename)}"
    )


def _draft_storage_path(
    *,
    company_id: UUID,
    checksum_sha256: str,
    original_filename: str,
) -> str:
    return (
        f"demo/company-kb/{company_id}/"
        f"{checksum_sha256[:12]}-{_safe_filename(original_filename)}"
    )


def _read_valid_draft_pdf(path: Path) -> bytes:
    if not path.exists():
        raise CompanyKbRegistrationError(f"PDF file does not exist: {path}")
    if not path.is_file():
        raise CompanyKbRegistrationError(f"PDF path is not a file: {path}")
    if path.suffix.lower() != ".pdf":
        raise CompanyKbRegistrationError(f"Input file is not a PDF: {path}")

    pdf_bytes = path.read_bytes()
    if not pdf_bytes.startswith(b"%PDF-"):
        raise CompanyKbRegistrationError(f"Input file is not a PDF: {path}")
    return pdf_bytes


def _normalize_draft_attachment_type(value: str) -> str:
    normalized = value.strip().lower()
    if normalized not in _DRAFT_ATTACHMENT_TYPES:
        raise CompanyKbRegistrationError(
            "attachment_type must be one of: "
            f"{', '.join(sorted(_DRAFT_ATTACHMENT_TYPES))}."
        )
    return normalized


def _safe_filename(filename: str) -> str:
    path = Path(filename)
    suffix = path.suffix.lower()
    stem = _slug(path.stem).lower() or "document"
    return f"{stem}{suffix}"


def _split_text(text: str, *, max_chunk_chars: int) -> list[str]:
    normalized = _normalize_text(text)
    if len(normalized) <= max_chunk_chars:
        return [normalized] if normalized else []
    parts: list[str] = []
    current: list[str] = []
    current_len = 0
    for sentence in re.split(r"(?<=[.!?])\s+", normalized):
        if not sentence:
            continue
        if current and current_len + len(sentence) + 1 > max_chunk_chars:
            parts.append(" ".join(current))
            current = []
            current_len = 0
        current.append(sentence)
        current_len += len(sentence) + 1
    if current:
        parts.append(" ".join(current))
    return parts


def _summary_for_document_type(
    document_type: CompanyKbDocumentType,
    text: str,
) -> str:
    excerpt = _excerpt(text)
    labels = {
        CompanyKbDocumentType.CERTIFICATION: "Certification evidence",
        CompanyKbDocumentType.CASE_STUDY: "Case study evidence",
        CompanyKbDocumentType.CV_PROFILE: "Anonymized CV/profile evidence",
        CompanyKbDocumentType.CAPABILITY_STATEMENT: "Capability evidence",
        CompanyKbDocumentType.POLICY_PROCESS: "Policy/process evidence",
        CompanyKbDocumentType.FINANCIAL_PRICING: "Financial/pricing evidence",
        CompanyKbDocumentType.LEGAL_INSURANCE: "Legal/insurance evidence",
    }
    return f"{labels[document_type]}: {excerpt}"


def _anonymize_cv_text(text: str) -> str:
    text = re.sub(r"(?i)\bname\s*:\s*.*?(?=\bemail\s*:)", "", text)
    text = re.sub(r"(?i)\bname\s*:\s*[^\n.]+", "", text)
    lines: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if re.match(r"(?i)^name\s*:", stripped):
            continue
        lines.append(stripped)
    out = "\n".join(lines)
    out = re.sub(
        r"(?i)\bemail\s*:\s*[\w.+-]+@[\w-]+(?:\.[\w-]+)+",
        "",
        out,
    )
    out = re.sub(r"[\w.+-]+@[\w-]+(?:\.[\w-]+)+", "", out)
    out = re.sub(
        r"(?i)\b(phone|mobile|tel|address)\s*:\s*\+?\d[\d\s().-]{6,}\d",
        "",
        out,
    )
    out = re.sub(r"\+?\d[\d\s().-]{6,}\d", "", out)
    return _normalize_text(out)


def _excerpt(text: str, max_len: int = 260) -> str:
    normalized = _normalize_text(text)
    if len(normalized) <= max_len:
        return normalized
    truncated = normalized[:max_len].rstrip()
    if " " in truncated:
        return truncated.rsplit(" ", 1)[0]
    return truncated


_PDF_ARTIFACT_LETTERS = "A-Za-zÅÄÖåäöÉéÈèÜü"


def _repair_pdf_text_extraction_artifacts(text: str) -> str:
    """Repair common bad ToUnicode mappings from exported Office PDFs.

    Some Calibri-based PDFs map ligature-like glyphs to punctuation or digits
    even though the visual PDF is correct. Keep this intentionally narrow so
    normal technical symbols and numbers are not rewritten.
    """

    repaired = text
    repaired = re.sub(r"(?<!\S)>ll\b", "till", repaired)
    repaired = re.sub(
        rf"\b>(?=[{_PDF_ARTIFACT_LETTERS}])",
        "ti",
        repaired,
    )
    repaired = re.sub(
        rf"(?<=[{_PDF_ARTIFACT_LETTERS}])>(?=[{_PDF_ARTIFACT_LETTERS}])",
        "ti",
        repaired,
    )
    repaired = re.sub(
        rf"(?<=[{_PDF_ARTIFACT_LETTERS}])0(?=[{_PDF_ARTIFACT_LETTERS}]|[,.!?;:]|\s|$)",
        "tt",
        repaired,
    )
    return _repair_truncated_tt_words(repaired)


def _repair_truncated_tt_words(text: str) -> str:
    words = sorted(
        set(re.findall(rf"\b[{_PDF_ARTIFACT_LETTERS}]{{5,}}tt\b", text)),
        key=len,
        reverse=True,
    )
    repaired = text
    for word in words:
        prefix = word[:-2]
        if len(prefix) < 4:
            continue
        repaired = re.sub(
            rf"\b{re.escape(prefix)}\.(?=\s|$)",
            word,
            repaired,
        )
    return repaired


def _normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _decode_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def _mapping(value: Any) -> Mapping[str, Any]:
    return value if isinstance(value, Mapping) else {}


def _normalize_uuid(value: Any, field_name: str) -> UUID:
    if isinstance(value, UUID):
        return value
    try:
        return UUID(str(value))
    except (TypeError, ValueError) as exc:
        raise CompanyKbError(f"{field_name} must be a valid UUID.") from exc


def _slug(value: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "-", value).strip("-").upper()
    return slug or "UNKNOWN"


__all__ = [
    "CompanyKbDocumentSummary",
    "CompanyKbDocumentType",
    "CompanyKbError",
    "CompanyKbFactExtractor",
    "CompanyKbIngestionResult",
    "CompanyKbPdfRegistrationResult",
    "CompanyKbRegistrationError",
    "CompanyKbRegistrationResult",
    "CompanyKbUploadFile",
    "EvidenceExcerptRef",
    "ExtractedCompanyKbFact",
    "ExtractedCompanyKbFacts",
    "RuleBasedCompanyKbFactExtractor",
    "build_company_kb_chunks",
    "delete_company_kb_document",
    "ingest_company_kb_document",
    "list_company_kb_documents",
    "list_company_kb_evidence",
    "register_company_kb_pdf",
    "register_company_kb_documents",
]
